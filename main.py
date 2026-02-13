import argparse
import os
import tempfile

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches


def pdf_to_images(pdf_path, out_dir, dpi=200, fmt="png"):
    doc = fitz.open(pdf_path)
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)

    paths = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        out_path = os.path.join(out_dir, f"page_{i + 1:04d}.{fmt}")
        pix.save(out_path)
        paths.append(out_path)

    doc.close()
    return paths


def images_to_pptx_full_bleed(image_paths, pptx_path):
    prs = Presentation()

    # 16:9（ワイド）固定
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    w, h = prs.slide_width, prs.slide_height

    for img in image_paths:
        slide = prs.slides.add_slide(blank)
        # 16:9前提で全面貼り（比率計算なし）
        slide.shapes.add_picture(img, 0, 0, width=w, height=h)

    prs.save(pptx_path)


def main():
    ap = argparse.ArgumentParser(description="Convert 16:9 PDF pages to images and build a 16:9 PPTX.")
    ap.add_argument("pdf", help="Input PDF path")
    ap.add_argument("pptx", help="Output PPTX path")
    ap.add_argument("--dpi", type=int, default=200, help="Render DPI (default: 200)")
    ap.add_argument("--format", choices=["png", "jpg"], default="png", help="Image format (default: png)")
    ap.add_argument("--keep-images", action="store_true", help="Keep rendered images next to PPTX")
    args = ap.parse_args()

    pdf_path = os.path.abspath(args.pdf)
    pptx_path = os.path.abspath(args.pptx)

    if not os.path.exists(pdf_path):
        raise FileNotFoundError(pdf_path)

    if args.keep_images:
        out_dir = os.path.join(os.path.dirname(pptx_path), "pdf_pages")
        os.makedirs(out_dir, exist_ok=True)
        imgs = pdf_to_images(pdf_path, out_dir, dpi=args.dpi, fmt=args.format)
        images_to_pptx_full_bleed(imgs, pptx_path)
        print(f"Saved: {pptx_path}")
        print(f"Images kept in: {out_dir}")
    else:
        with tempfile.TemporaryDirectory() as tmpdir:
            imgs = pdf_to_images(pdf_path, tmpdir, dpi=args.dpi, fmt=args.format)
            images_to_pptx_full_bleed(imgs, pptx_path)
        print(f"Saved: {pptx_path}")


if __name__ == "__main__":
    main()
